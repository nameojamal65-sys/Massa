#!/usr/bin/env python3
import os, datetime
from pathlib import Path
import openai
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

# -------------------
# إعداد المشروع والشركة
# -------------------
PROJECT_DIR = Path.home() / "sovereign_production"
OUTPUT_DIR = PROJECT_DIR / "enterprise_output"
OUTPUT_DIR.mkdir(exist_ok=True)

# -------------------
# جمع ملفات المشروع
# -------------------
files = [str(f) for f in PROJECT_DIR.rglob("*.py")]
print(f"🚀 عدد الملفات المكتشفة: {len(files)}")

# -------------------
# تحسين الملفات ودمج قدرات إدارة الشركة
# -------------------
def ai_optimize(files_list):
    optimized_files = {}
    for file_path in files_list:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            prompt = f"""
            طور هذا الكود ليصبح نظام إدارة مقاولات وشركة ذكي (Enterprise AI Construction & Company Management System) 
            يشمل كل الأقسام، التقارير اليومية/الأسبوعية/الشهرية، صلاحيات متعددة للمستخدمين، إنتاج PowerPoint وPresentations تلقائي.
            حافظ على وظائف الكود الأصلي وأضف هذه الميزات.
            الكود الحالي:
            {content}
            """
            response = openai.ChatCompletion.create(
                model="gpt-4-turbo",
                messages=[{"role":"user","content":prompt}],
                temperature=0.3,
                max_tokens=4000
            )
            optimized_files[file_path] = response.choices[0].message["content"]
        except Exception as e:
            print(f"⚠️ خطأ في معالجة {file_path}: {e}")
    return optimized_files

optimized = ai_optimize(files)

# -------------------
# حفظ الملفات المحسنة
# -------------------
for path, content in optimized.items():
    save_path = OUTPUT_DIR / Path(path).relative_to(PROJECT_DIR)
    save_path.parent.mkdir(parents=True, exist_ok=True)
    with open(save_path, "w", encoding="utf-8") as f:
        f.write(content)

# -------------------
# إعداد تقارير Excel وPowerPoint
# -------------------
def generate_report_excel(df, report_name):
    file_path = OUTPUT_DIR / f"{report_name}.xlsx"
    df.to_excel(file_path, index=False)
    print(f"📄 تقرير Excel جاهز: {file_path}")

def generate_ppt_from_df(df, report_name):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Report: {report_name}"
    for i, row in df.iterrows():
        slide.shapes.add_textbox(Inches(1), Inches(1+i*0.5), Inches(8), Inches(0.5)).text = str(row.to_dict())
    ppt_path = OUTPUT_DIR / f"{report_name}.pptx"
    prs.save(ppt_path)
    print(f"📊 Presentation جاهز: {ppt_path}")

# -------------------
# مثال على بيانات المشاريع
# -------------------
data = {
    "Project": ["Project A", "Project B"],
    "Status": ["On Track", "Delayed"],
    "Progress %": [45, 30],
}
df = pd.DataFrame(data)
generate_report_excel(df, "weekly_report")
generate_ppt_from_df(df, "weekly_report")

# -------------------
# نظام صلاحيات المستخدمين
# -------------------
users = {
    "admin": {"role": "Project Manager", "access": "full"},
    "user_hr": {"role": "HR", "access": ["HR", "Training"]},
    "user_qa": {"role": "QA/QC", "access": ["QA/QC"]},
    "owner": {"role": "Owner", "access": ["Summary Reports"]},
    "consultant": {"role": "Consultant", "access": ["Project Status"]},
}

print("✅ صلاحيات المستخدمين جاهزة")

print(f"✅ العملية اكتملت! الملفات المحسنة في {OUTPUT_DIR}")

